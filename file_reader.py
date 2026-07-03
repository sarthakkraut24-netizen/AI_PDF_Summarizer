# ==========================================================
# FILE_READER.PY
# AI Document Assistant
# PyMuPDF + EasyOCR Version
# Compatible with Python 3.13+
# ==========================================================

import os
import io

import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image

import fitz  # PyMuPDF

from utils import clean_text
from ocr import (
    extract_text_from_scanned_pdf,
    needs_ocr,
    extract_text_from_image
)

from table_reader import read_tables
from image_reader import process_image

# ==========================================================
# SUPPORTED FILE TYPES
# ==========================================================

SUPPORTED_FILES = [
    ".pdf",
    ".docx",
    ".txt",
    ".pptx",
    ".xlsx",
    ".xls",
    ".png",
    ".jpg",
    ".jpeg",
    ".bmp",
    ".tiff",
    ".webp"
]

# ==========================================================
# GET FILE EXTENSION
# ==========================================================

def get_extension(uploaded_file):

    return os.path.splitext(uploaded_file.name)[1].lower()

# ==========================================================
# READ NORMAL PDF (TEXT BASED)
# ==========================================================

def read_pdf(file):

    try:

        file.seek(0)

        reader = PdfReader(file)

        text = ""

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:

                text += page_text + "\n"

        return clean_text(text)

    except Exception as e:

        raise Exception(f"Unable to read PDF:\n{e}")

# ==========================================================
# SCANNED PDF READER (OCR)
# ==========================================================

def read_scanned_pdf(file):

    try:

        file.seek(0)

        text = extract_text_from_scanned_pdf(file)

        return clean_text(text)

    except Exception as e:

        raise Exception(f"OCR Failed:\n{e}")

# ==========================================================
# SMART PDF HANDLER
# ==========================================================

def extract_pdf(file):

    text = read_pdf(file)

    if needs_ocr(text):

        print("Scanned PDF detected... switching to OCR")

        text = read_scanned_pdf(file)

    return clean_text(text)
# ==========================================================
# READ DOCX FILE
# ==========================================================

def read_docx(file):

    try:

        file.seek(0)

        document = Document(file)

        text = []

        for para in document.paragraphs:

            if para.text.strip():

                text.append(para.text.strip())

        return clean_text("\n".join(text))

    except Exception as e:

        raise Exception(f"Unable to read DOCX:\n{e}")


# ==========================================================
# READ TXT FILE
# ==========================================================

def read_txt(file):

    encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]

    for enc in encodings:

        try:

            file.seek(0)

            text = file.read().decode(enc)

            return clean_text(text)

        except:

            continue

    raise Exception("Unable to decode TXT file")


# ==========================================================
# READ PPTX FILE
# ==========================================================

def read_pptx(file):

    try:

        file.seek(0)

        presentation = Presentation(file)

        text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():

                        text.append(shape.text.strip())

        return clean_text("\n".join(text))

    except Exception as e:

        raise Exception(f"Unable to read PPTX:\n{e}")


# ==========================================================
# READ EXCEL FILE
# ==========================================================

def read_excel(file):

    try:

        file.seek(0)

        excel_file = pd.ExcelFile(file)

        all_text = []

        for sheet in excel_file.sheet_names:

            df = pd.read_excel(excel_file, sheet_name=sheet, dtype=str)

            df.fillna("", inplace=True)

            all_text.append(f"=== {sheet} ===")

            all_text.append(df.to_string(index=False))

        return clean_text("\n\n".join(all_text))

    except Exception as e:

        raise Exception(f"Unable to read Excel:\n{e}")
 # ==========================================================
# READ IMAGE FILE
# ==========================================================

def read_image(file):

    try:

        result = process_image(file)

        return clean_text(result["text"])

    except Exception as e:

        raise Exception(f"Unable to read image:\n{e}")


# ==========================================================
# EXTRACT TABLES FROM PDF
# ==========================================================

def extract_pdf_tables(file):

    try:

        file.seek(0)

        result = read_tables(file)

        return result.get("text", "")

    except Exception:

        return ""


# ==========================================================
# MAIN FILE READER (CORE FUNCTION)
# ==========================================================

def read_file(uploaded_file):

    if uploaded_file is None:

        return ""

    extension = get_extension(uploaded_file)

    # -------------------------
    # PDF FILE
    # -------------------------

    if extension == ".pdf":

        text = extract_pdf(uploaded_file)

        table_text = extract_pdf_tables(uploaded_file)

        if table_text.strip():

            text += "\n\n=== TABLE DATA ===\n"

            text += table_text

        return clean_text(text)

    # -------------------------
    # DOCX FILE
    # -------------------------

    elif extension == ".docx":

        return read_docx(uploaded_file)

    # -------------------------
    # TXT FILE
    # -------------------------

    elif extension == ".txt":

        return read_txt(uploaded_file)

    # -------------------------
    # PPTX FILE
    # -------------------------

    elif extension == ".pptx":

        return read_pptx(uploaded_file)

    # -------------------------
    # EXCEL FILE
    # -------------------------

    elif extension in [".xlsx", ".xls"]:

        return read_excel(uploaded_file)

    # -------------------------
    # IMAGE FILES
    # -------------------------

    elif extension in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:

        return read_image(uploaded_file)

    # -------------------------
    # UNSUPPORTED
    # -------------------------

    else:

        raise ValueError(f"Unsupported file type: {extension}")
# ==========================================================
# FILE INFO
# ==========================================================

def file_info(uploaded_file):

    if uploaded_file is None:

        return {}

    return {

        "name": uploaded_file.name,

        "size_kb": round(uploaded_file.size / 1024, 2),

        "extension": get_extension(uploaded_file)

    }


# ==========================================================
# SUPPORTED FILE TYPES
# ==========================================================

def supported_files():

    return SUPPORTED_FILES


# ==========================================================
# QUICK VALIDATION
# ==========================================================

def is_supported_file(filename):

    ext = os.path.splitext(filename)[1].lower()

    return ext in SUPPORTED_FILES


# ==========================================================
# TEST BLOCK
# ==========================================================

if __name__ == "__main__":

    print("File Reader Loaded Successfully")

    print("\nSupported File Types:")

    for f in SUPPORTED_FILES:

        print(" -", f)

    print("\nModule Ready 🚀")
           